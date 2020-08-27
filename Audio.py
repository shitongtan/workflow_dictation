#speech.py
import pyaudio
import wave
import requests
import json
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.dml.color import RGBColor, ColorFormat
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

API_ENDPOINT = 'https://api.wit.ai/speech'
ACCESS_TOKEN = '2BG6V4SKNCEK4JBNZ2QHG7BV7GZYGSDL'
action_words=[]

def record_audio(RECORD_SECONDS, WAVE_OUTPUT_FILENAME):
    #--------- SETTING PARAMS FOR OUR AUDIO FILE ------------#
    FORMAT = pyaudio.paInt16    # format of wave
    CHANNELS = 1                # no. of audio channels
    RATE = 44100                # frame rate
    CHUNK = 1024                # frames per audio sample
    #--------------------------------------------------------#

    # creating PyAudio object
    audio = pyaudio.PyAudio()

    # open a new stream for microphone
    # It creates a PortAudio Stream Wrapper class object
    stream = audio.open(format=FORMAT,channels=CHANNELS,
                        rate=RATE, input=True,
                        frames_per_buffer=CHUNK)


    #----------------- start of recording -------------------#
    print("Listening...")

    # list to save all audio frames
    frames = []

    for i in range(int(RATE / CHUNK * RECORD_SECONDS)):
        # read audio stream from microphone
        data = stream.read(CHUNK)
        # append audio data to frames list
        frames.append(data)

    #------------------ end of recording --------------------#   
    print("Finished recording.")

    stream.stop_stream()    # stop the stream object
    stream.close()          # close the stream object
    audio.terminate()       # terminate PortAudio

    #------------------ saving audio ------------------------#

    # create wave file object
    waveFile = wave.open(WAVE_OUTPUT_FILENAME, 'wb')

    # settings for wave file object
    waveFile.setnchannels(CHANNELS)
    waveFile.setsampwidth(audio.get_sample_size(FORMAT))
    waveFile.setframerate(RATE)
    waveFile.writeframes(b''.join(frames))

    # closing the wave file object
    waveFile.close()


def read_audio(WAVE_FILENAME):
    # function to read audio(wav) file
    with open(WAVE_FILENAME, 'rb') as f:
        audio = f.read()
    return audio

def RecognizeSpeech(AUDIO_FILENAME, num_seconds = 5):

    # record audio of specified length in specified audio file
    record_audio(num_seconds, AUDIO_FILENAME)

    # reading audio
    audio = read_audio(AUDIO_FILENAME)

    # WIT.AI HERE
    # defining headers for HTTP request
    headers = {'authorization': 'Bearer ' + ACCESS_TOKEN,
               'Content-Type': 'audio/wav'}
    
    #Send the request as post request and the audio as data
    resp = requests.post(API_ENDPOINT, headers = headers,
                             data = audio)

    #Get the text
    data = json.loads(resp.content)
    if 'wit$message_body:message_body' in data['entities']:
        print(data['entities']['wit$message_body:message_body'][0]['value'])
        word=(data['entities']['wit$message_body:message_body'][0]['value'])
        action_words.append(word)
    else:
        print("Empty")
        return "Empty"

#if __name__ == "__main__":
def run_audio():
    text=""
    while text!="Empty":
        text = RecognizeSpeech('myspeech.wav', 10)
    prs = Presentation()
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    
    shapes.title.text = 'Your flowchart'
    
    left = Inches(0.93)  # 0.93" centers this overall set of shapes
    top = Inches(3.0)
    width = Inches(1.75)
    height = Inches(1.0)
    
    shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
    text_frame = shape.text_frame
    text_frame.auto_size= MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = action_words[0]
    font = run.font
    font.color.rgb = RGBColor(0, 0, 0)
    font.size = Pt(12)

    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)


    left = left + width - Inches(0.4)
    width = Inches(2.0)  # chevrons need more width for visual balance
    
    for n in range(1,len(action_words)):
        shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
        text_frame = shape.text_frame
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = action_words[n]
        font = run.font
        font.color.rgb = RGBColor(0, 0, 0)
        font.size = Pt(12)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        left = left + width - Inches(0.4)

    out_file = io.BytesIO()
    prs.save(out_file)
    out_file.seek(0)
    return out_file